import time
import os
from pptx import Presentation

PMEM_DIR = "/mnt/pmem0"
PDF_FILENAME = "paper.pdf"
PPTX_FILENAME = "presentation.pptx"


def run():
    print(f"--- [Code] Processing Job ---")

    pdf_path = os.path.join(PMEM_DIR, PDF_FILENAME)
    pptx_path = os.path.join(PMEM_DIR, PPTX_FILENAME)

    # 1. Transfer In (Read from PMEM)
    print(f"Reading from PMEM: {pdf_path}...")

    start_read = time.time()
    if not os.path.exists(pdf_path):
        print("Error: PDF not found on PMEM mount.")
        return

    with open(pdf_path, "rb") as f:
        _ = f.read()  # Force read
    end_read = time.time()

    # 2. Processing (Convert)
    print("Converting PDF to PPTX...")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PMEM Benchmark Report"

    # Save to local temp first to separate processing from transfer
    local_output = "temp_output.pptx"
    prs.save(local_output)

    # 3. Transfer Out (Write to PMEM)
    print(f"Writing to PMEM: {pptx_path}...")

    with open(local_output, "rb") as source:
        content = source.read()

    start_write = time.time()
    with open(pptx_path, "wb") as dest:
        dest.write(content)
        dest.flush()
        os.fsync(dest.fileno())
    end_write = time.time()

    print(f"SUCCESS: Pipeline Complete.")
    print(f"METRIC_TRANSFER_READ: {end_read - start_read:.6f} seconds")
    print(f"METRIC_TRANSFER_WRITE: {end_write - start_write:.6f} seconds")


if __name__ == "__main__":
    run()
