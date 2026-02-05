import time
import os
from pptx import Presentation

NAS_DIR = "/mnt/nas"
PDF_FILENAME = "paper.pdf"
PPTX_FILENAME = "presentation.pptx"


def run():
    print(f"--- [Code] Processing Job ---")

    pdf_path = os.path.join(NAS_DIR, PDF_FILENAME)
    pptx_path = os.path.join(NAS_DIR, PPTX_FILENAME)

    # 1. Transfer In (Read from NAS)
    print(f"Reading from NAS: {pdf_path}...")

    start_read = time.time()
    if not os.path.exists(pdf_path):
        print("Error: PDF not found on NAS mount.")
        return

    with open(pdf_path, "rb") as f:
        _ = f.read()  # Force read over network
    end_read = time.time()

    # 2. Processing (Convert)
    print("Converting PDF to PPTX...")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "NAS Benchmark Report"

    # Save to local temp first to separate processing from transfer
    local_output = "temp_output.pptx"
    prs.save(local_output)

    # 3. Transfer Out (Write to NAS)
    print(f"Writing to NAS: {pptx_path}...")

    with open(local_output, "rb") as source:
        content = source.read()

    start_write = time.time()
    with open(pptx_path, "wb") as dest:
        dest.write(content)
        dest.flush()
        os.fsync(dest.fileno())  # Force network flush
    end_write = time.time()

    print(f"SUCCESS: Pipeline Complete.")
    print(f"METRIC_TRANSFER_READ: {end_read - start_read:.6f} seconds")
    print(f"METRIC_TRANSFER_WRITE: {end_write - start_write:.6f} seconds")


if __name__ == "__main__":
    run()
